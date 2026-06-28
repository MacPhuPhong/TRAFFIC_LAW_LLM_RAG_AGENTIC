# -*- coding: utf-8 -*-
"""
Dựng slide báo cáo ĐATN (30 phút) — end-to-end — đồng bộ với BaoCao_MacPhuPhong_3.pdf.
SV: Mạc Phú Phong (106210059) · GVHD: TS. Trần Thị Minh Hạnh · ĐH Bách Khoa Đà Nẵng.
Nguồn sự thật: doc/BaoCao_MacPhuPhong_3.pdf. Biểu đồ: research/results/figures|plots/*.png.
Chạy: /home/pphong/venv/LLM_Agentic/bin/python3 doc/slides_v7/build_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image

# ---------- Đường dẫn ----------
ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
FIG = os.path.join(ROOT, "research", "results", "figures")
PLOT = os.path.join(ROOT, "research", "results", "plots")
OUT = os.path.join(os.path.dirname(__file__), "Slide_BaoCao_MacPhuPhong_30phut.pptx")

def fig(name):
    p = os.path.join(FIG, name)
    if not os.path.exists(p):
        raise FileNotFoundError(p)
    return p

def plot(name):
    p = os.path.join(PLOT, name)
    if not os.path.exists(p):
        raise FileNotFoundError(p)
    return p

# ---------- Bảng màu ----------
NAVY   = RGBColor(0x1A, 0x23, 0x7E)   # offline / chủ đạo
BLUE   = RGBColor(0x15, 0x65, 0xC0)   # api / dense
TEAL   = RGBColor(0x00, 0x69, 0x5C)   # client
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
AMBER  = RGBColor(0xF5, 0x7F, 0x17)   # agent
ORANGE = RGBColor(0xE6, 0x5C, 0x00)   # core / retrieval
PINK   = RGBColor(0xC2, 0x18, 0x5B)   # store
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)   # ext
RED    = RGBColor(0xC6, 0x28, 0x28)   # HITL / cảnh báo
GREY   = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0xEE, 0xEE, 0xEE)
DGREY  = RGBColor(0x33, 0x33, 0x33)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x21, 0x21, 0x21)
CREAM  = RGBColor(0xF7, 0xF6, 0xF2)

FONT = "Arial"
EMU_IN = 914400
SW, SH = 13.333, 7.5   # inch (16:9)

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

_slide_no = 0  # đếm slide hiển thị (bìa không đánh số)

# =================================================================
# Helper cấp thấp
# =================================================================
def _tf(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    return tf

def _set_text(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return p

def textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = _tf(tb); tf.vertical_anchor = anchor
    _set_text(tf, text, size, color, bold, align, font)
    return tb

def box(slide, x, y, w, h, text, fill, fg=WHITE, size=12, bold=True,
        line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = _tf(sh); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    first = True
    for line_txt in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run(); r.text = line_txt
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = fg; f.name = FONT
    return sh

def arrow(slide, x1, y1, x2, y2, color=GREY, width=1.6, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(d)
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn

def chip(slide, x, y, text, fill, w=2.0, h=0.34, size=11):
    box(slide, x, y, w, h, text, fill, WHITE, size, True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)

# ---- Khung chung của 1 slide nội dung ----
def base(kicker=None, kicker_color=NAVY):
    global _slide_no
    s = prs.slides.add_slide(BLANK)
    # nền
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    # dải accent trên cùng
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    strip.fill.solid(); strip.fill.fore_color.rgb = kicker_color; strip.line.fill.background()
    strip.shadow.inherit = False
    _slide_no += 1
    # footer
    textbox(s, 0.4, 7.06, 9.5, 0.35, "ĐATN · Hỏi đáp luật GTĐB ứng dụng LLM kết hợp RAG · Mạc Phú Phong",
            10, GREY, False, PP_ALIGN.LEFT)
    textbox(s, 11.6, 7.06, 1.33, 0.35, str(_slide_no), 10, GREY, False, PP_ALIGN.RIGHT)
    if kicker:
        chip(s, 0.4, 0.28, kicker, kicker_color, w=max(1.4, 0.16*len(kicker)), h=0.32, size=11)
    return s

def title_on(s, title, y=0.72, color=NAVY, size=27):
    textbox(s, 0.4, y, 12.5, 1.0, title, size, color, True)

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def bullets_box(s, x, y, w, h, items, base_size=18, gap=6):
    """items: list of (level, text) hoặc str. level 0/1/2."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = _tf(tb); tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for it in items:
        if isinstance(it, tuple):
            lvl, txt = it
        else:
            lvl, txt = 0, it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        bullet = {0: "▸  ", 1: "–  ", 2: "·  "}.get(lvl, "•  ")
        # màu/đậm theo cấp
        size = base_size - (lvl * 2)
        # tách **bold** thô sơ: nếu text bắt đầu bằng '##' coi như heading nhỏ
        r = p.add_run()
        r.text = ("" if lvl == 0 and txt.startswith("§") else bullet) + txt.replace("§", "")
        f = r.font; f.size = Pt(size); f.name = FONT
        f.bold = (lvl == 0)
        f.color.rgb = INK if lvl == 0 else DGREY
    return tb

def fit(img_path, max_w, max_h):
    w, h = Image.open(img_path).size
    ar = w / h
    if max_w / max_h > ar:
        return ar * max_h, max_h
    return max_w, max_w / ar

def image_slide(kicker, kc, title, img_path, caption=None, side=None, notes_text=""):
    s = base(kicker, kc); title_on(s, title)
    if side:
        # ảnh trái, bullet phải
        iw, ih = fit(img_path, 7.7, 4.7)
        ix = 0.5 + (7.9 - iw) / 2
        s.shapes.add_picture(img_path, Inches(ix), Inches(1.85 + (4.9 - ih)/2),
                             width=Inches(iw), height=Inches(ih))
        bullets_box(s, 8.55, 1.9, 4.4, 4.8, side, base_size=16, gap=7)
    else:
        iw, ih = fit(img_path, 11.0, 4.85)
        ix = (SW - iw) / 2
        s.shapes.add_picture(img_path, Inches(ix), Inches(1.85),
                             width=Inches(iw), height=Inches(ih))
    if caption:
        textbox(s, 0.5, 6.6, 12.3, 0.4, caption, 12, GREY, False, PP_ALIGN.CENTER)
    if notes_text:
        notes(s, notes_text)
    return s

def table_slide(kicker, kc, title, headers, rows, notes_text="", col_w=None,
                hl_rows=None, font_size=14, note_line=None):
    s = base(kicker, kc); title_on(s, title)
    hl_rows = hl_rows or []
    nrows = len(rows) + 1; ncols = len(headers)
    top = 1.85
    height = min(4.7, 0.42 * nrows)
    gt = s.shapes.add_table(nrows, ncols, Inches(0.5), Inches(top),
                            Inches(12.3), Inches(height)).table
    if col_w:
        total = sum(col_w)
        for i, c in enumerate(col_w):
            gt.columns[i].width = Inches(12.3 * c / total)
    # header
    for j, htxt in enumerate(headers):
        cell = gt.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = kc
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = htxt
        r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = gt.cell(i, j)
            if i in hl_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else LGREY
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(font_size - 1); r.font.name = FONT
            r.font.bold = (i in hl_rows)
            r.font.color.rgb = INK
    if note_line:
        textbox(s, 0.5, top + height + 0.15, 12.3, 0.8, note_line, 13, GREY, False)
    if notes_text:
        notes(s, notes_text)
    return s

def bullets_slide(kicker, kc, title, items, notes_text="", subtitle=None, base_size=18):
    s = base(kicker, kc); title_on(s, title)
    y = 1.8
    if subtitle:
        textbox(s, 0.5, 1.62, 12.3, 0.5, subtitle, 15, kc, True)
        y = 2.2
    bullets_box(s, 0.6, y, 12.2, 6.5 - y, items, base_size=base_size, gap=8)
    if notes_text:
        notes(s, notes_text)
    return s

def two_col_slide(kicker, kc, title, left_title, left_items, right_title, right_items,
                  notes_text="", lc=NAVY, rc=ORANGE):
    s = base(kicker, kc); title_on(s, title)
    box(s, 0.5, 1.75, 6.0, 0.5, left_title, lc, WHITE, 15, True)
    box(s, 6.83, 1.75, 6.0, 0.5, right_title, rc, WHITE, 15, True)
    bullets_box(s, 0.6, 2.45, 5.85, 4.3, left_items, base_size=16, gap=7)
    bullets_box(s, 6.93, 2.45, 5.85, 4.3, right_items, base_size=16, gap=7)
    if notes_text:
        notes(s, notes_text)
    return s

# =================================================================
# SLIDE 1 — Bìa
# =================================================================
def slide_title():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.55))
    band.fill.solid(); band.fill.fore_color.rgb = WHITE; band.line.fill.background(); band.shadow.inherit = False
    textbox(s, 0.8, 0.62, 11.7, 0.45, "ĐẠI HỌC ĐÀ NẴNG  ·  TRƯỜNG ĐẠI HỌC BÁCH KHOA",
            14, RGBColor(0xC5, 0xCA, 0xE9), True, PP_ALIGN.CENTER)
    textbox(s, 0.8, 1.08, 11.7, 0.4, "KHOA ĐIỆN TỬ - VIỄN THÔNG",
            13, RGBColor(0xC5, 0xCA, 0xE9), True, PP_ALIGN.CENTER)
    textbox(s, 0.8, 1.62, 11.7, 0.55, "ĐỒ ÁN TỐT NGHIỆP",
            20, WHITE, True, PP_ALIGN.CENTER)
    textbox(s, 0.8, 2.12, 11.7, 0.4, "Ngành: Điện tử - Viễn thông  ·  Chuyên ngành: Hệ thống máy tính",
            12, RGBColor(0xC5, 0xCA, 0xE9), False, PP_ALIGN.CENTER)
    textbox(s, 0.6, 2.72, 12.1, 1.55,
            "XÂY DỰNG HỆ THỐNG HỎI ĐÁP LUẬT GIAO THÔNG ĐƯỜNG BỘ\n"
            "VIỆT NAM ỨNG DỤNG LLM KẾT HỢP RAG",
            26, NAVY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, 0.8, 4.4, 11.7, 0.5,
            "Truy xuất lai · Phân đoạn phân cấp pháp luật · LangGraph · Con người trong vòng lặp",
            14, ORANGE, True, PP_ALIGN.CENTER)
    textbox(s, 0.8, 5.35, 11.7, 1.3,
            "Sinh viên thực hiện:  Mạc Phú Phong  —  106210059  —  Lớp 21DTCLC1\n"
            "Giảng viên hướng dẫn:  TS. Trần Thị Minh Hạnh",
            16, WHITE, False, PP_ALIGN.CENTER)
    textbox(s, 0.8, 6.8, 11.7, 0.4, "Đà Nẵng, tháng 6 năm 2026", 13, RGBColor(0xC5,0xCA,0xE9), False, PP_ALIGN.CENTER)
    notes(s, "(~0:30 | cộng dồn 0:30)\n"
             "Kính chào quý thầy cô trong hội đồng. Em là Mạc Phú Phong, xin trình bày đồ án tốt "
             "nghiệp: Xây dựng hệ thống hỏi đáp luật giao thông đường bộ Việt Nam ứng dụng LLM kết "
             "hợp RAG, dưới sự hướng dẫn của cô Trần Thị Minh Hạnh. Báo cáo gồm 4 phần: bối cảnh & "
             "cơ sở lý thuyết, thiết kế hệ thống, thực nghiệm đánh giá, và triển khai – kết luận. "
             "Em xin phép bắt đầu.")

# =================================================================
# SLIDE 2 — Agenda
# =================================================================
def slide_agenda():
    s = base("Nội dung", NAVY); title_on(s, "Nội dung trình bày")
    items = [
        ("1", "Bối cảnh, bài toán & cơ sở lý thuyết", "Vì sao cần RAG cho pháp luật giao thông", NAVY),
        ("2", "Thiết kế & xây dựng hệ thống", "Pipeline ngoại tuyến + trực tuyến, agent LangGraph", ORANGE),
        ("3", "Thực nghiệm & đánh giá (10 thực nghiệm)", "Bằng chứng định lượng cho từng quyết định thiết kế", GREEN),
        ("4", "Triển khai, minh hoạ & kết luận", "Production $0/tháng, demo, hạn chế & hướng phát triển", PURPLE),
    ]
    y = 1.95
    for no, t, d, c in items:
        box(s, 0.8, y, 0.75, 0.95, no, c, WHITE, 30, True)
        box(s, 1.75, y, 10.7, 0.95, "", LGREY, INK, 12, False, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, 2.0, y+0.12, 10.2, 0.45, t, 19, c, True)
        textbox(s, 2.0, y+0.55, 10.2, 0.35, d, 13, GREY, False)
        y += 1.15
    notes(s, "(~0:20 | 0:50)\n"
             "Phần một, em nói về bối cảnh và lý do chọn kiến trúc RAG. Phần hai là thiết kế hệ "
             "thống end-to-end. Phần ba là mười thực nghiệm đánh giá. Phần cuối là triển khai thực "
             "tế và kết luận. Trọng tâm của báo cáo nằm ở phần hai và phần ba.")

# =================================================================
# SLIDE 3 — Đặt vấn đề
# =================================================================
def slide_problem():
    s = base("Phần 1 · Bối cảnh", NAVY)
    title_on(s, "Đặt vấn đề & động lực nghiên cứu")
    box(s, 0.5, 1.75, 6.0, 0.5, "Bối cảnh: cải cách pháp luật 2024–2026", NAVY, WHITE, 15, True)
    bullets_box(s, 0.6, 2.45, 5.9, 4.2, [
        (0, "Hàng loạt văn bản mới hiệu lực gần đồng thời"),
        (1, "Luật TTATGT 36/2024 · Luật Đường bộ 35/2024"),
        (1, "NĐ 168/2024: tăng mạnh mức phạt + trừ điểm GPLX (cơ chế hoàn toàn mới)"),
        (0, "Nhu cầu tra cứu chính xác tới Điều/Khoản/Điểm tăng vọt"),
    ], base_size=16, gap=9)
    box(s, 6.83, 1.75, 6.0, 0.5, "Hạn chế của giải pháp hiện tại", ORANGE, WHITE, 15, True)
    bullets_box(s, 6.93, 2.45, 5.9, 4.2, [
        (0, "Cổng pháp luật: phải tự biết tên VB & số điều; không hỏi ngôn ngữ tự nhiên"),
        (0, "LLM thương mại: chưa cập nhật corpus 2024–2026"),
        (1, "→ Ảo giác pháp lý (hallucination): trích sai nguồn, nhầm mức phạt"),
        (0, "Rủi ro nghiêm trọng trong lĩnh vực pháp lý"),
    ], base_size=16, gap=9)
    box(s, 0.5, 6.5, 12.3, 0.6,
        "Giải pháp: RAG hướng tác tử — truy xuất văn bản gốc rồi mới sinh, kiểm soát nguồn trích dẫn, hạn chế ảo giác",
        RED, WHITE, 14, True)
    notes(s, "(~1:10 | 2:00)\n"
             "Pháp luật giao thông Việt Nam vừa trải qua đợt cải cách lớn nhất nhiều năm: Luật 35 và "
             "36/2024, đặc biệt Nghị định 168/2024 tăng mạnh mức phạt và lần đầu áp dụng trừ điểm "
             "giấy phép lái xe. Người dân, CSGT và đơn vị tư vấn đều cần tra cứu chính xác tới từng "
             "điểm, khoản, điều. Nhưng cổng pháp luật chính thức buộc người dùng tự biết tên văn bản "
             "và số điều, không hiểu câu hỏi tự nhiên; còn các mô hình ngôn ngữ lớn thương mại chưa "
             "cập nhật corpus 2024–2026 nên hay bịa nguồn, nhầm mức phạt — gọi là ảo giác pháp lý, "
             "rất nguy hiểm. RAG giải quyết bằng cách truy xuất văn bản gốc đưa vào ngữ cảnh trước "
             "khi sinh câu trả lời, nhờ đó kiểm soát được nguồn trích dẫn.")

# =================================================================
# SLIDE 4 — Mục tiêu, phạm vi, đóng góp
# =================================================================
def slide_goals():
    s = base("Phần 1 · Bối cảnh", NAVY)
    title_on(s, "Mục tiêu · Phạm vi · Đóng góp")
    box(s, 0.5, 1.7, 6.0, 0.46, "Mục tiêu & phạm vi", NAVY, WHITE, 15, True)
    bullets_box(s, 0.6, 2.25, 5.9, 4.5, [
        (0, "Trợ lý hỏi đáp tiếng Việt, trả lời theo văn bản gốc, trích dẫn tới Điểm/Khoản/Điều"),
        (0, "Corpus: 24 văn bản còn hiệu lực (06/2026)"),
        (1, "2 Luật · 5 Nghị định · 17 Thông tư"),
        (1, "NĐ 81/2026 đường sắt: bao quát giao cắt đường bộ–đường sắt"),
        (0, "Không cần GPU cho vận hành (trừ reranker tuỳ chọn)"),
        (0, "Ngoài phạm vi: đường thuỷ, hàng không; tư vấn cá nhân hoá; phụ lục mẫu biểu"),
    ], base_size=15, gap=7)
    box(s, 6.83, 1.7, 6.0, 0.46, "3 đóng góp chính", ORANGE, WHITE, 15, True)
    bullets_box(s, 6.93, 2.25, 5.9, 4.5, [
        (0, "Thuật toán phân đoạn phân cấp Điều→Khoản→Điểm"),
        (1, "5 cải tiến; điểm thiếu đoạn 77,5%→11,7%; Cit-R ×2,4"),
        (0, "Quy trình truy xuất lai đa truy vấn đặc thù VB pháp luật"),
        (1, "Mở rộng QV · RRF · sibling · giải tham chiếu chéo + bộ đánh giá 40 câu/8 danh mục (corpus 2024–2026 mới)"),
        (0, "Kiến trúc kiểm soát ảo giác pháp lý theo chiều sâu (3 lớp)"),
        (1, "Lệnh hướng dẫn (×8,7) · lọc trích dẫn · HITL"),
    ], base_size=15, gap=7)
    notes(s, "(~1:00 | 3:00)\n"
             "Mục tiêu là xây trợ lý hỏi đáp tiếng Việt trả lời đúng theo văn bản gốc với trích dẫn "
             "tới cấp điểm, khoản, điều. Phạm vi gồm 24 văn bản còn hiệu lực: 2 luật, 5 nghị định, "
             "17 thông tư — trong đó có bổ sung Nghị định 81/2026 về đường sắt để bao quát các quy "
             "định giao cắt đường bộ – đường sắt; ngoài phạm vi là pháp luật đường thuỷ, hàng không "
             "và tư vấn cá nhân hoá. Đề tài có ba đóng góp chính: thứ nhất là thuật toán phân đoạn "
             "phân cấp theo cấu trúc Điều–Khoản–Điểm với năm cải tiến; thứ hai là quy trình truy xuất "
             "lai đa truy vấn đặc thù cho văn bản pháp luật kèm bộ đánh giá 40 câu trên corpus giao "
             "thông 2024–2026 chưa từng được xử lý; thứ ba là kiến trúc kiểm soát ảo giác pháp lý "
             "theo chiều sâu gồm ba lớp — lệnh hướng dẫn ràng buộc trích dẫn, lọc trích dẫn sau sinh, "
             "và con người trong vòng lặp.")

# =================================================================
# SLIDE 5 — RAG là gì
# =================================================================
def slide_rag():
    s = base("Phần 1 · Cơ sở lý thuyết", NAVY)
    title_on(s, "RAG — Sinh văn bản tăng cường truy xuất")
    # flow ngang
    y = 2.6; h = 0.95; w = 1.95
    xs = [0.5, 2.65, 4.8, 6.95, 9.1]
    labels = [("Truy vấn\nngười dùng", BLUE), ("Mã hoá\ntruy vấn (E5)", BLUE),
              ("Tìm kiếm vector\ntop-k đoạn", ORANGE), ("Ghép ngữ cảnh\n+ truy vấn", ORANGE),
              ("LLM sinh\ncâu trả lời", AMBER)]
    for (lab, c), x in zip(labels, xs):
        box(s, x, y, w, h, lab, c, WHITE, 13, True)
    for i in range(len(xs)-1):
        arrow(s, xs[i]+w, y+h/2, xs[i+1], y+h/2, GREY, 1.8)
    box(s, 11.2, y, 1.7, h, "Câu trả lời\n+ trích dẫn", GREEN, WHITE, 12, True)
    arrow(s, xs[4]+w, y+h/2, 11.2, y+h/2, GREY, 1.8)
    box(s, 4.8, 1.55, 2.1, 0.7, "Kho tài liệu\nkiểm soát (Qdrant)", PINK, WHITE, 12, True)
    arrow(s, 5.85, 2.25, 5.77, y, PINK, 1.8)
    bullets_box(s, 0.6, 3.95, 12.2, 2.5, [
        (0, "Không nhồi kiến thức vào tham số → cập nhật văn bản không cần huấn luyện lại"),
        (0, "Kiểm soát nguồn trích dẫn → hạn chế ảo giác (đặc biệt quan trọng với pháp luật)"),
        (0, "Thách thức đặc thù: cấu trúc phân cấp Điều→Khoản→Điểm; truy vấn ngắn/mơ hồ; câu đa ý định"),
    ], base_size=16, gap=9)
    textbox(s, 0.5, 6.55, 12.3, 0.4, "Hình 3 — Kiến trúc tổng quan RAG (Lewis và cộng sự, 2020)",
            12, GREY, False, PP_ALIGN.CENTER)
    notes(s, "(~1:10 | 4:10)\n"
             "RAG hoạt động qua các bước: mã hoá truy vấn thành vector, tìm các đoạn văn gần nhất "
             "trong kho tài liệu được kiểm soát, ghép chúng vào ngữ cảnh rồi để LLM sinh câu trả lời "
             "kèm trích dẫn. Ưu điểm: không nhồi kiến thức vào tham số nên cập nhật văn bản mà không "
             "phải huấn luyện lại, và kiểm soát được nguồn nên hạn chế ảo giác. Với pháp luật Việt "
             "Nam có ba thách thức đặc thù mà đề tài phải xử lý: cấu trúc phân cấp Điều–Khoản–Điểm, "
             "truy vấn ngắn mơ hồ, và câu hỏi đa ý định.")

# =================================================================
# SLIDE 6 — Nền tảng kỹ thuật
# =================================================================
def slide_foundations():
    s = base("Phần 1 · Cơ sở lý thuyết", NAVY)
    title_on(s, "Nền tảng kỹ thuật của truy xuất lai")
    cols = [
        ("Truy xuất DÀY ĐẶC", BLUE, [
            "Embedding E5-base đa ngôn ngữ (XLM-R)",
            "768 chiều · gộp trung bình + chuẩn hoá L2",
            "Bắt ngữ nghĩa: 'vượt đèn đỏ' ≈ 'không chấp hành hiệu lệnh đèn'",
        ]),
        ("Truy xuất THƯA (BM25)", ORANGE, [
            "Tần suất từ, không cần học",
            "k₁=1,5 · b=0,75 · tách từ pyvi",
            "Mạnh với mã chính xác: 'Điều 7', 'NĐ 168/2024'",
        ]),
        ("Hợp nhất RRF (k=60)", GREEN, [
            "Gộp nhiều danh sách xếp hạng",
            "Không cần chuẩn hoá điểm giữa hai hệ",
            "Ổn định cho corpus < 10.000 tài liệu",
        ]),
    ]
    x = 0.5
    for t, c, items in cols:
        box(s, x, 1.8, 4.0, 0.5, t, c, WHITE, 14, True)
        bullets_box(s, x+0.08, 2.45, 3.95, 3.6, [(0 if i == 0 else 1, it) for i, it in enumerate(items)],
                    base_size=14, gap=8)
        x += 4.17
    box(s, 0.5, 6.2, 12.3, 0.65,
        "Dày đặc + Thưa bổ sung cho nhau → truy xuất lai vượt trội từng phương pháp đơn lẻ (Luan và cộng sự)",
        NAVY, WHITE, 14, True)
    notes(s, "(~1:20 | 5:30)\n"
             "Truy xuất lai kết hợp hai nhánh bổ sung cho nhau. Nhánh dày đặc dùng embedding E5-base "
             "đa ngôn ngữ 768 chiều, bắt được tương đồng ngữ nghĩa — ví dụ 'vượt đèn đỏ' và 'không "
             "chấp hành hiệu lệnh đèn tín hiệu' tuy khác từ nhưng cùng nghĩa. Nhánh thưa dùng BM25 dựa "
             "trên tần suất từ, rất mạnh khi truy vấn chứa mã chính xác như 'Điều 7' hay 'Nghị định "
             "168'. Hai danh sách được hợp nhất bằng RRF với k bằng 60 — phương pháp gộp không cần "
             "chuẩn hoá điểm, ổn định cho kho dưới mười nghìn tài liệu. Kết hợp hai nhánh cho kết quả "
             "vượt trội từng nhánh riêng lẻ.")

# =================================================================
# SLIDE 7 — Agentic + LangGraph + HITL
# =================================================================
def slide_agentic():
    s = base("Phần 1 · Cơ sở lý thuyết", NAVY)
    title_on(s, "Hướng tác tử · LangGraph · Con người trong vòng lặp")
    box(s, 0.5, 1.75, 6.0, 0.46, "LangGraph — máy trạng thái cho LLM", AMBER, WHITE, 14, True)
    bullets_box(s, 0.6, 2.3, 5.9, 4.5, [
        (0, "Nút (Node): xử lý + cập nhật trạng thái"),
        (0, "Cạnh có điều kiện: định tuyến đa ý định"),
        (0, "Trạng thái (AgentState): ngữ cảnh xuyên suốt"),
        (0, "Checkpointer: Sqlite (dev) / Postgres (prod)"),
        (1, "Lưu bền → hội thoại đa lượt, không mất state"),
        (0, "interrupt_before: tạm dừng chờ phê duyệt"),
    ], base_size=15, gap=7)
    box(s, 6.83, 1.75, 6.0, 0.46, "HITL & khoảng trống nghiên cứu", RED, WHITE, 14, True)
    bullets_box(s, 6.93, 2.3, 5.9, 4.5, [
        (0, "Con người trong vòng lặp: duyệt thông tin rủi ro cao trước khi trả về"),
        (1, "Áp dụng cho luồng tra cứu web (Tavily)"),
        (0, "Khoảng trống: RAG pháp luật VN mới tập trung Bộ luật Lao động/Dân sự"),
        (1, "Chưa xử lý corpus giao thông 2024–2026, cấu trúc Điều→Khoản→Điểm, trừ điểm GPLX"),
    ], base_size=15, gap=8)
    notes(s, "(~1:00 | 6:30)\n"
             "Hệ thống hướng tác tử có khả năng lập kế hoạch và ra quyết định qua nhiều bước. Đề tài "
             "dùng LangGraph — thư viện mô hình hoá ứng dụng LLM thành máy trạng thái với các nút, "
             "cạnh có điều kiện, trạng thái dùng chung, và bộ lưu điểm kiểm tra. Điểm quan trọng: "
             "checkpointer lưu bền vào SQLite khi phát triển và Postgres khi production, nên hỗ trợ "
             "hội thoại nhiều lượt và cho phép tạm dừng đồ thị chờ phê duyệt mà không mất trạng thái. "
             "Cơ chế con người trong vòng lặp dùng để kiểm duyệt thông tin tra cứu web trước khi trả "
             "về. Khoảng trống nghiên cứu: các hệ RAG pháp luật tiếng Việt trước đây chủ yếu làm Bộ "
             "luật Lao động hay Dân sự, chưa ai xử lý corpus giao thông 2024–2026 với cấu trúc phân "
             "cấp phức tạp và cơ chế trừ điểm mới — đây chính là chỗ đề tài lấp đầy.")

# =================================================================
# SLIDE 8 — Kiến trúc tổng thể
# =================================================================
def slide_arch():
    s = base("Phần 2 · Thiết kế", ORANGE)
    title_on(s, "Kiến trúc tổng thể: hai pipeline tách biệt")
    # offline row
    box(s, 0.5, 1.7, 2.6, 0.42, "NGOẠI TUYẾN (offline)", NAVY, WHITE, 13, True, align=PP_ALIGN.LEFT)
    oy = 2.25; oh = 0.8; ow = 2.05
    oxs = [0.5, 2.75, 5.0, 7.25]
    off = [("24 văn bản\nPDF", NAVY), ("Làm sạch +\nChunking v6", NAVY),
           ("Nhúng E5-base\n768 chiều", NAVY), ("", PINK)]
    for (lab, c), x in zip(off, oxs):
        if lab:
            box(s, x, oy, ow, oh, lab, c, WHITE, 12, True)
    for i in range(2):
        arrow(s, oxs[i]+ow, oy+oh/2, oxs[i+1], oy+oh/2, GREY, 1.6)
    arrow(s, oxs[2]+ow, oy+oh/2, 9.55, oy+oh/2, GREY, 1.6)
    qd = box(s, 9.55, oy-0.05, 3.2, 0.9, "Qdrant — Traffic_Law_Hybrid\n4.423 đoạn × 768d · HNSW\n+ BM25 cache (pyvi)",
             PINK, WHITE, 11, True)
    # online row
    box(s, 0.5, 3.55, 2.6, 0.42, "TRỰC TUYẾN (online)", AMBER, WHITE, 13, True, align=PP_ALIGN.LEFT)
    ny = 4.1; nh = 0.8
    box(s, 0.5, ny, 1.7, nh, "Người\ndùng", TEAL, WHITE, 12, True)
    box(s, 2.45, ny, 2.0, nh, "analyzer\ný định + mở rộng QV", AMBER, WHITE, 11, True)
    box(s, 4.7, ny, 1.6, nh, "Định tuyến\n4 nhánh", AMBER, WHITE, 11, True)
    arrow(s, 2.2, ny+nh/2, 2.45, ny+nh/2, GREY, 1.6)
    arrow(s, 4.45, ny+nh/2, 4.7, ny+nh/2, GREY, 1.6)
    branches = [("legal_rag — truy xuất lai + sinh P2", ORANGE, 4.35),
                ("chit_chat", GREY, 5.05),
                ("web_search → web_finalize (HITL)", RED, 5.75),
                ("out_of_scope", GREY, 6.45)]
    for lab, c, yy in branches:
        box(s, 6.6, yy, 4.4, 0.5, lab, c, WHITE, 11, True)
        arrow(s, 6.3, ny+nh/2, 6.6, yy+0.25, GREY, 1.2)
    box(s, 11.3, 4.9, 1.5, 1.0, "Câu trả lời\n+ trích dẫn\n(SSE)", GREEN, WHITE, 11, True)
    for _, _, yy in branches:
        arrow(s, 11.0, yy+0.25, 11.3, 5.4, GREY, 1.0)
    arrow(s, 11.15, oy+0.9, 6.9, 4.35, PINK, 1.4, dashed=True)  # Qdrant -> legal_rag
    notes(s, "(~1:30 | 8:00)\n"
             "Hệ thống tách thành hai pipeline. Pipeline ngoại tuyến chỉ chạy khi thêm văn bản: 24 "
             "PDF được làm sạch, chunking phân cấp v6, nhúng bằng E5-base 768 chiều rồi nạp vào "
             "Qdrant — collection Traffic_Law_Hybrid 4.423 đoạn, kèm cache BM25. Pipeline trực tuyến "
             "xử lý mỗi câu hỏi: nút analyzer phân tích ý định và mở rộng truy vấn, router định tuyến "
             "về một trong bốn nhánh — hỏi đáp pháp luật, hội thoại thường, tra cứu web có phê duyệt, "
             "và ngoài phạm vi. Nhánh legal_rag truy vấn Qdrant rồi sinh câu trả lời, phát về giao "
             "diện theo từng token qua SSE. Tách hai pipeline giúp cập nhật corpus không ảnh hưởng "
             "tính sẵn sàng phục vụ. (Tương ứng Hình 4–5 trong báo cáo.)")

# =================================================================
# SLIDE 9 — Tech stack
# =================================================================
def slide_stack():
    rows = [
        ["Điều phối agent", "LangGraph 0.2.x", "State machine, HITL interrupt_before"],
        ["LLM sinh", "Gemini 3.1 Flash Lite", "1M context · $0,075/$0,30 per 1M token"],
        ["Embedding", "multilingual-e5-base", "768d · R@10 cao nhất trong 4 model"],
        ["Vector DB", "Qdrant 1.17", "Payload filter · HTTP API · HNSW"],
        ["Sparse", "rank_bm25 + pyvi", "Tách từ ghép tiếng Việt"],
        ["Backend", "FastAPI + uvicorn", "Async · SSE streaming"],
        ["Frontend", "Next.js 14 + Tailwind", "Auth gate · citation panel"],
        ["Judge eval", "OpenAI gpt-4o-mini", "RAGAS · judge độc lập"],
        ["Tra cứu web", "Tavily", "HITL · free tier"],
        ["Quan sát", "LangSmith", "Phân tích chi phí node-level"],
    ]
    table_slide("Phần 2 · Thiết kế", ORANGE, "Công nghệ sử dụng (tech stack)",
                ["Tầng", "Thư viện / Dịch vụ", "Lý do chọn"], rows,
                col_w=[2.6, 4.0, 6.0], font_size=14,
                notes_text="(~0:30 | 8:30)\n"
                "Toàn bộ tech stack: LangGraph điều phối agent; Gemini 3.1 Flash Lite sinh câu trả "
                "lời với cửa sổ ngữ cảnh 1 triệu token và chi phí rất thấp; embedding e5-base; vector "
                "DB Qdrant; BM25 với tách từ pyvi cho tiếng Việt; backend FastAPI streaming SSE; "
                "frontend Next.js 14; đánh giá độc lập bằng gpt-4o-mini qua RAGAS; tra cứu web Tavily; "
                "và LangSmith để phân tích chi phí từng nút. Em sẽ đi sâu từng thành phần.")

# =================================================================
# SLIDE 10 — Làm sạch PDF
# =================================================================
def slide_clean():
    s = base("Phần 2 · Ngoại tuyến", ORANGE)
    title_on(s, "Bước 1 — Chuẩn hoá văn bản pháp luật (PDF → Markdown)")
    box(s, 0.5, 1.75, 6.0, 0.46, "3 script làm sạch chuyên biệt (pdfplumber)", NAVY, WHITE, 14, True)
    bullets_box(s, 0.6, 2.3, 5.9, 3.0, [
        (0, "clean_luat — Luật: phần/chương/điều, header mỗi trang"),
        (0, "clean_nghidinh — NĐ: bảng mức phạt → Markdown"),
        (0, "clean_thongtu — TT: phụ lục, tiêu đề lồng nhau"),
        (1, "Lọc header/footer, nối câu bị ngắt trang, bảo toàn bảng"),
    ], base_size=15, gap=8)
    box(s, 6.83, 1.75, 6.0, 0.46, "Lọc văn bản hết hiệu lực (trường status)", ORANGE, WHITE, 14, True)
    bullets_box(s, 6.93, 2.3, 5.9, 3.0, [
        (0, "active → đưa vào index"),
        (0, "repealed → bỏ qua (vd Luật GTĐB 2008)"),
        (0, "partially_repealed → bỏ mảng đường bộ (vd NĐ 100/2019 bị NĐ 168 thay)"),
    ], base_size=15, gap=10)
    box(s, 0.5, 5.55, 12.3, 0.95,
        "Vì sao quan trọng: văn bản pháp luật là PDF layout phức tạp (bảng phạt, ngắt dòng giữa câu).\n"
        "Làm sạch sai → chunking & trích dẫn sai theo. Lọc status đảm bảo không trả về điều luật đã bị bãi bỏ.",
        CREAM, INK, 14, False)
    notes(s, "(~1:00 | 9:30)\n"
             "Bước đầu của pipeline ngoại tuyến là chuẩn hoá. Văn bản pháp luật phát hành dưới dạng "
             "PDF có layout phức tạp: bảng mức phạt, header/footer lặp mỗi trang, câu bị ngắt giữa "
             "chừng do sang trang. Em viết ba script làm sạch riêng cho ba loại — luật, nghị định, "
             "thông tư — dùng pdfplumber, lọc header/footer, nối câu bị ngắt, và bảo toàn bảng dưới "
             "dạng Markdown. Một điểm quan trọng là trường status: văn bản còn hiệu lực mới được "
             "index; văn bản bị bãi bỏ như Luật Giao thông đường bộ 2008 thì bỏ qua; văn bản bị bãi "
             "bỏ một phần như Nghị định 100/2019 thì loại mảng đường bộ. Nhờ vậy hệ thống không bao "
             "giờ trả về điều luật đã hết hiệu lực.")

# =================================================================
# SLIDE 11 — Chunking phân cấp v6
# =================================================================
def slide_chunk():
    s = base("Phần 2 · Ngoại tuyến", ORANGE)
    title_on(s, "Bước 2 — Chunking phân cấp v6 (đóng góp cốt lõi)")
    # cây phân cấp
    box(s, 0.5, 1.8, 3.0, 0.7, "L1 · ĐIỀU\n≤ 1.000 token → 1 chunk", NAVY, WHITE, 12, True)
    box(s, 0.5, 2.65, 3.0, 0.7, "L2 · KHOẢN\n≤ 500 token → 1 chunk", BLUE, WHITE, 12, True)
    box(s, 0.5, 3.5, 3.0, 0.7, "L3 · ĐIỂM\nchunk độc lập + làm giàu", ORANGE, WHITE, 12, True)
    arrow(s, 2.0, 2.5, 2.0, 2.65, GREY, 1.4)
    arrow(s, 2.0, 3.35, 2.0, 3.5, GREY, 1.4)
    box(s, 0.5, 4.5, 3.0, 1.7, "Phân bố 4.423 chunk\n\nL1 Điều: 486 (11%)\nL2 Khoản: 1.282 (29%)\nL3 Điểm: 2.655 (60%)",
        CREAM, INK, 13, True)
    box(s, 3.8, 1.8, 9.0, 0.46, "5 cải tiến của v6 — giải bài toán Điểm thiếu chunk", ORANGE, WHITE, 14, True)
    bullets_box(s, 3.9, 2.35, 8.9, 2.6, [
        (0, "Giữ Điểm ngắn (≥5 token), không lọc bỏ"),
        (0, "Ngưỡng strict 80 token cho nghị định"),
        (0, "NĐ có ≥2 Điểm → luôn tách L3"),
        (0, "Regex bắt cả 'đ)' và Điểm sau dấu chấm phẩy"),
        (0, "Làm giàu: chunk Điểm tự chứa phần mở đầu Khoản ('Phạt tiền X–Y đồng đối với…')"),
    ], base_size=15, gap=7)
    box(s, 3.8, 5.15, 9.0, 1.05,
        "Kết quả: tỉ lệ Điểm pháp lý của NĐ 168 KHÔNG có chunk độc lập giảm 77,5% → 11,7%.\n"
        "Biên chunk TRÙNG biên trích dẫn pháp lý → lấy đúng chunk là tự có đúng Khoản để trích dẫn.",
        RED, WHITE, 14, True)
    notes(s, "(~1:40 | 11:10)\n"
             "Đây là đóng góp cốt lõi của đề tài. Bộ chunking phân cấp v6 tách văn bản theo đúng cấu "
             "trúc pháp lý ba cấp: nếu một Điều ngắn dưới 1.000 token thì giữ nguyên thành một chunk; "
             "nếu dài thì tách theo Khoản; nếu Khoản vẫn dài hoặc là nghị định có nhiều Điểm thì tách "
             "tiếp xuống từng Điểm. Phiên bản v6 có năm cải tiến quan trọng: giữ lại các Điểm ngắn, "
             "hạ ngưỡng cho nghị định, ép tách Điểm khi nghị định có từ hai điểm trở lên, sửa regex "
             "để bắt cả điểm 'đ' và điểm sau dấu chấm phẩy, và đặc biệt là làm giàu — mỗi chunk Điểm "
             "tự mang theo phần mở đầu của Khoản, ví dụ câu 'Phạt tiền từ X đến Y đồng đối với…'. Nhờ "
             "vậy tỉ lệ điểm pháp lý của Nghị định 168 không có chunk riêng giảm từ 77,5% xuống còn "
             "11,7%. Cái hay nằm ở chỗ: biên của chunk trùng với biên của đơn vị trích dẫn — nên khi "
             "retriever lấy đúng chunk là tự động có đúng khoản, điểm để trích dẫn. Tổng cộng 4.423 "
             "chunk, trong đó 60% là cấp Điểm.")

# =================================================================
# SLIDE 12 — Embedding + Index
# =================================================================
def slide_index():
    s = base("Phần 2 · Ngoại tuyến", ORANGE)
    title_on(s, "Bước 3 — Nhúng vector & đánh chỉ mục vào Qdrant")
    two = [
        (0, "E5-base mã hoá 4.423 chunk, batch 64"),
        (1, "Bắt buộc prefix 'passage: ' khi nhúng đoạn, 'query: ' khi nhúng truy vấn"),
        (1, "Dùng sai prefix → giảm ~30% recall (vấn đề kỹ thuật #1)"),
        (0, "Collection Traffic_Law_Hybrid: 768d, Cosine, HNSW m=16/ef=100"),
        (0, "3 payload index KEYWORD: doc_id, status, topic → lọc nhanh tại query time"),
        (0, "Mỗi chunk gắn metadata đầy đủ: chunk_id, doc_id, dieu, khoan, diem, level…"),
        (0, "BM25 cache: tokenize pyvi rồi lưu JSONL → khởi động lại < 0,5s (không rebuild)"),
        (0, "Tổng dung lượng sau index: ~26 MB"),
    ]
    bullets_box(s, 0.6, 1.85, 12.2, 4.9, two, base_size=17, gap=10)
    notes(s, "(~1:00 | 12:10)\n"
             "Bước ba là nhúng và đánh chỉ mục. Mô hình e5-base mã hoá toàn bộ 4.423 chunk. Một yêu "
             "cầu bắt buộc của họ E5 là thêm tiền tố: 'passage:' khi nhúng đoạn văn và 'query:' khi "
             "nhúng truy vấn — đây cũng là vấn đề kỹ thuật đầu tiên em gặp, dùng sai prefix làm giảm "
             "khoảng 30% recall. Collection Qdrant cấu hình 768 chiều, khoảng cách cosine, HNSW. Em "
             "tạo ba payload index trên doc_id, status và topic để lọc nhanh — ví dụ chỉ lấy status "
             "active ngay tại tầng cơ sở dữ liệu. Mỗi chunk mang metadata đầy đủ tới cấp điểm khoản "
             "điều. Cache BM25 được tokenize sẵn bằng pyvi và lưu JSONL nên khởi động lại dưới nửa "
             "giây thay vì rebuild. Toàn bộ index chỉ khoảng 26 MB.")

# =================================================================
# SLIDE 13 — Truy xuất lai 5 bước
# =================================================================
def slide_retrieval():
    s = base("Phần 2 · Trực tuyến", ORANGE)
    title_on(s, "Truy xuất lai — pipeline 5 bước")
    steps = [
        ("1 · Mở rộng truy vấn", "analyzer sinh 1–5 biến thể; giải tham chiếu hội thoại", BLUE),
        ("2 · Dày đặc + Thưa", "E5 (lọc status=active) + BM25 Okapi → mỗi nhánh top-2k", ORANGE),
        ("3 · Hợp nhất RRF (k=60)", "gộp mọi danh sách → top-30 không trùng", GREEN),
        ("4 · Làm giàu + giải tham chiếu chéo", "sibling ±2 khoản · regex 'Điều X Khoản Y' · diversity cap", AMBER),
        ("5 · (tuỳ chọn) Xếp hạng lại", "cross-encoder bge-reranker → giữ top-10", PINK),
    ]
    y = 1.85
    for t, d, c in steps:
        box(s, 0.6, y, 4.2, 0.78, t, c, WHITE, 14, True)
        textbox(s, 5.05, y+0.06, 7.8, 0.7, d, 14, INK, False, MSO_ANCHOR.MIDDLE)
        if y < 5.0:
            arrow(s, 2.7, y+0.78, 2.7, y+0.92, GREY, 1.4)
        y += 0.92
    box(s, 0.6, 6.45, 12.25, 0.5,
        "Sibling & cross-ref đặc trị NĐ 168: mức phạt tiền và mức trừ điểm ở hai khoản khác nhau, liên kết qua tham chiếu",
        NAVY, WHITE, 13, True)
    notes(s, "(~1:40 | 13:50)\n"
             "Khi có câu hỏi, pipeline truy xuất chạy năm bước. Bước một, analyzer mở rộng truy vấn "
             "thành một đến năm biến thể và giải tham chiếu theo lịch sử hội thoại. Bước hai, chạy "
             "song song nhánh dày đặc E5 — đã lọc status active ngay tại Qdrant — và nhánh thưa BM25, "
             "mỗi nhánh lấy top 2k. Bước ba, hợp nhất tất cả danh sách bằng RRF ra top 30 không trùng. "
             "Bước bốn là phần đặc thù pháp luật: sibling enrichment kéo thêm các khoản lân cận cùng "
             "Điều, và giải tham chiếu chéo dùng regex tìm các 'Điều X Khoản Y' được nhắc tới rồi kéo "
             "luôn chunk đó vào — điều này đặc trị Nghị định 168, nơi mức phạt tiền và mức trừ điểm "
             "nằm ở hai khoản khác nhau nhưng liên kết qua tham chiếu. Cuối cùng áp diversity cap rồi "
             "cắt top 10. Bước năm là reranker tuỳ chọn, mặc định tắt — em sẽ giải thích ở phần đánh "
             "giá.")

# =================================================================
# SLIDE 14 — Prompt P2
# =================================================================
def slide_prompt():
    s = base("Phần 2 · Trực tuyến", ORANGE)
    title_on(s, "Sinh văn bản — Lệnh hướng dẫn (14 quy tắc) & lọc trích dẫn")
    box(s, 0.5, 1.75, 6.0, 0.46, "Lệnh hướng dẫn đầy đủ quy tắc (P2)", NAVY, WHITE, 14, True)
    bullets_box(s, 0.6, 2.3, 5.9, 4.3, [
        (0, "Chỉ dùng NGỮ CẢNH, không suy luận ngoài → chống ảo giác"),
        (0, "Mọi mức phạt/điểm phải kèm [doc_id, Điều, Khoản, Điểm]"),
        (0, "Không đủ thông tin → từ chối an toàn"),
        (0, "Phân biệt phương tiện: Đ6 ô tô · Đ7 mô tô · Đ8 chuyên dùng"),
        (0, "Cross-reference: ghép mức phạt tiền với mức trừ điểm"),
        (0, "Câu đa ý → tách bullet, in đậm số tiền/điểm"),
    ], base_size=14, gap=6)
    box(s, 6.83, 1.75, 6.0, 0.46, "Đầu ra có cấu trúc + Citation Sanitation", ORANGE, WHITE, 14, True)
    bullets_box(s, 6.93, 2.3, 5.9, 4.3, [
        (0, "Pydantic GeneratedAnswer: answer + sources[]"),
        (1, "Gemini trả JSON đã validate, không parse chuỗi"),
        (0, "Lọc trích dẫn (hậu xử lý): so từng citation với metadata 10 chunk thực sự trong prompt"),
        (1, "Citation không khớp (vd Điều 99 bịa) → loại bỏ trước khi tới người dùng"),
        (0, "→ Lớp phòng thủ cuối chống ảo giác trích dẫn"),
    ], base_size=14, gap=8)
    notes(s, "(~1:00 | 14:50)\n"
             "Sau khi có ngữ cảnh, bộ sinh dùng lệnh hướng dẫn đầy đủ quy tắc (biến thể P2 trong biểu "
             "đồ) gồm phần khai báo vai trò và 14 quy tắc bắt "
             "buộc. Vài quy tắc quan trọng: chỉ dùng ngữ cảnh được cung cấp; mọi mức phạt phải kèm "
             "trích dẫn tới điểm khoản điều; nếu thiếu thông tin thì từ chối thay vì bịa; phân biệt "
             "đúng điều theo loại phương tiện — Điều 6 ô tô, Điều 7 mô tô; và ghép mức phạt tiền với "
             "mức trừ điểm qua tham chiếu chéo. Gemini trả về JSON có cấu trúc theo Pydantic nên "
             "không phải parse chuỗi thủ công. Cuối cùng là lớp lọc trích dẫn: hệ thống so từng "
             "citation mà LLM trả về với metadata của đúng 10 chunk đã đưa vào prompt; nếu LLM bịa ra "
             "một điều không có trong ngữ cảnh thì citation đó bị loại trước khi tới người dùng — đây "
             "là lớp phòng thủ cuối cùng chống ảo giác trích dẫn.")

# =================================================================
# SLIDE 15 — LangGraph 6 nút + HITL
# =================================================================
def slide_graph():
    s = base("Phần 2 · Agent", ORANGE)
    title_on(s, "Agent LangGraph — 6 nút, định tuyến & HITL")
    box(s, 0.5, 1.85, 1.7, 0.7, "BẮT ĐẦU", TEAL, WHITE, 12, True)
    an = box(s, 2.5, 1.85, 2.3, 0.7, "analyzer\ný định·loại xe·QV", AMBER, WHITE, 12, True)
    rt = box(s, 5.1, 1.85, 1.9, 0.7, "định tuyến\ntheo danh mục", AMBER, WHITE, 12, True)
    arrow(s, 2.2, 2.2, 2.5, 2.2, GREY, 1.5)
    arrow(s, 4.8, 2.2, 5.1, 2.2, GREY, 1.5)
    nodes = [("legal_rag — truy xuất lai + P2", ORANGE, 1.55),
             ("chit_chat — template tĩnh", GREY, 2.45),
             ("out_of_scope — từ chối", GREY, 3.35),
             ("web_search (Tavily)", BLUE, 4.25)]
    for lab, c, yy in nodes:
        box(s, 7.6, yy, 4.0, 0.62, lab, c, WHITE, 12, True)
        arrow(s, 7.0, 2.2, 7.6, yy+0.31, GREY, 1.1)
    box(s, 11.75, 2.35, 1.1, 0.7, "KẾT\nTHÚC", TEAL, WHITE, 12, True)
    # web_finalize HITL — ngay dưới web_search, mũi tên không cắt nút khác
    box(s, 7.6, 5.15, 4.0, 0.62, "web_finalize — tạm dừng chờ duyệt (HITL)", RED, WHITE, 11, True)
    arrow(s, 9.6, 4.87, 9.6, 5.15, RED, 1.4)
    # các nút kết thúc → END (legal_rag, chit_chat, out_of_scope, web_finalize)
    for yy in (1.55, 2.45, 3.35):
        arrow(s, 11.6, yy+0.31, 11.78, 2.7, GREY, 1.0)
    arrow(s, 11.6, 5.46, 11.95, 3.05, GREY, 1.0)
    # HITL strip
    box(s, 0.5, 6.0, 12.35, 0.95,
        "HITL: web_search thu ≤3 snippet → interrupt_before['web_finalize'] → /pending → Admin Console\n"
        "Approve (resume web_finalize, gắn nhãn cảnh báo tra cứu internet) hoặc Reject (xoá thread). Timeout 10 phút.",
        CREAM, INK, 13, False)
    notes(s, "(~1:20 | 16:10)\n"
             "Toàn bộ luồng trực tuyến được điều phối bởi đồ thị LangGraph sáu nút, điểm vào là "
             "analyzer. Analyzer phân tích ý định, loại phương tiện, mở rộng truy vấn; router định "
             "tuyến về bốn nhánh. Nhánh legal_rag chạy truy xuất lai và sinh P2 rồi kết thúc. "
             "chit_chat và out_of_scope dùng template tĩnh, không tốn token LLM. Nhánh web_search gọi "
             "Tavily lấy tối đa ba đoạn trích rồi đồ thị tạm dừng tại web_finalize nhờ interrupt_before "
             "— đây là cơ chế con người trong vòng lặp. Quản trị viên vào trang admin, xem snippet, "
             "rồi Approve để đồ thị tiếp tục tổng hợp câu trả lời có gắn nhãn cảnh báo tra cứu "
             "internet, hoặc Reject để xoá. Có timeout 10 phút phía client. Nhờ checkpointer, trạng "
             "thái không mất khi tạm dừng, và mỗi phiên có thread_id riêng đảm bảo cách ly nhiều "
             "người dùng. (Đồ thị: Hình 7; luồng HITL: Hình 8.)")

# =================================================================
# SLIDE 16 — Thiết kế thực nghiệm
# =================================================================
def slide_exp_design():
    s = base("Phần 3 · Đánh giá", GREEN)
    title_on(s, "Thiết kế thực nghiệm — 40 câu, 8 danh mục, 10 thực nghiệm")
    box(s, 0.5, 1.75, 6.0, 0.46, "Bộ đánh giá 40 câu / 8 danh mục", GREEN, WHITE, 14, True)
    bullets_box(s, 0.6, 2.3, 5.9, 4.3, [
        (0, "5 câu mỗi danh mục, xây thủ công từ tình huống thực"),
        (1, "Phạt đơn nghĩa · Đa ý định · Tham chiếu chéo"),
        (1, "Thủ tục · Ngoài phạm vi · Ngắn, đa nghĩa"),
        (1, "Đa hành vi · Ràng buộc phương tiện"),
        (0, "Mỗi câu có đáp án chuẩn + trích dẫn chuẩn (cấp Khoản)"),
        (0, "5 câu ngoài phạm vi lọc khỏi đo truy xuất → giữ 35 câu pháp lý"),
    ], base_size=14, gap=7)
    box(s, 6.83, 1.75, 6.0, 0.46, "Bộ metric đa chiều", ORANGE, WHITE, 14, True)
    bullets_box(s, 6.93, 2.3, 5.9, 4.3, [
        (0, "Truy xuất: Recall@k, MRR, nDCG@10 (cấp Khoản & cấp tài liệu)"),
        (0, "Trích dẫn: Citation P / R / F1 (metric cốt lõi)"),
        (0, "Câu trả lời: Token-F1, ROUGE-L, Refusal rate"),
        (0, "Kiến trúc: Category Accuracy của router"),
        (0, "RAGAS (gpt-4o-mini): faithfulness, context precision/recall"),
        (0, "Vận hành: cost/câu, latency, error rate (LangSmith)"),
    ], base_size=14, gap=6)
    box(s, 0.5, 6.55, 12.3, 0.42,
        "10 thực nghiệm đối chứng theo thứ tự từ dưới lên (phân đoạn → nhúng → DB → lệnh HD → truy xuất → kiến trúc) — đổi 1 biến/lần",
        NAVY, WHITE, 12, True)
    notes(s, "(~1:00 | 17:10)\n"
             "Để kiểm chứng từng quyết định, em xây bộ đánh giá 40 câu hỏi chia tám danh mục thực "
             "tiễn, mỗi danh mục năm câu — từ câu phạt đơn giản, đa ý định, tham chiếu chéo, thủ tục, "
             "tới câu ngoài phạm vi, câu ngắn mơ hồ, hành động kép và phân biệt loại phương tiện. Mỗi "
             "câu có đáp án chuẩn và tập trích dẫn chuẩn ở cấp khoản. Năm câu ngoài phạm vi được lọc "
             "khỏi các thực nghiệm đo truy xuất, giữ lại 35 câu pháp lý. Bộ metric đa chiều: recall, "
             "MRR, nDCG cho truy xuất; citation precision/recall/F1 cho trích dẫn — đây là metric cốt "
             "lõi; token-F1 và ROUGE cho câu trả lời; category accuracy cho router; RAGAS với "
             "gpt-4o-mini làm trọng tài độc lập; và chi phí, độ trễ, tỉ lệ lỗi từ LangSmith. Mười "
             "thực nghiệm chạy theo thứ tự bottom-up, mỗi lần chỉ đổi một biến để tránh nhiễu.")

# =================================================================
# SLIDE 17 — RQ2 Chunking
# =================================================================
def slide_rq2():
    image_slide("Phần 3 · Mục 3.3", GREEN,
                "Phân đoạn phân cấp vs cố định 512 — phát hiện then chốt",
                fig("rq2_chunking_answer.png"),
                caption="Hình 10 — Đầu–cuối: F1, ROUGE-L, Citation-Recall, tỷ lệ từ chối (phân cấp vs cố định 512)",
                side=[
                    (0, "Citation-Recall (cấp Khoản):"),
                    (1, "0,410 vs 0,171 → ×2,4"),
                    (0, "Cit-R cấp tài liệu:"),
                    (1, "0,814 vs 0,457"),
                    (0, "Refusal: 0,20 vs 0,37"),
                    (0, "F1 token gần hoà (0,327 ≈ 0,325)"),
                    (1, "fixed-512 viết dài, trùng n-gram nhưng KHÔNG cite đúng"),
                    (0, "→ Biên chunk = biên trích dẫn là chìa khoá"),
                ],
                notes_text="(~1:10 | 18:20)\n"
                "Thực nghiệm phân đoạn (Mục 3.3) so sánh phân đoạn phân cấp với cắt cố định 512 token. Kết quả quan trọng nhất: "
                "Citation-Recall ở cấp khoản đạt 0,410 so với 0,171 — gấp 2,4 lần. Lý do như em đã "
                "nói: biên chunk trùng biên trích dẫn pháp lý nên lấy đúng chunk là có đúng khoản. "
                "Đáng chú ý, F1 token gần như hoà nhau, vì fixed-512 tạo chunk dài chứa nhiều n-gram "
                "trùng đáp án nhưng lại không trích dẫn đúng được. Đây chính là lý do trong domain "
                "pháp lý, Citation-Recall mới là thước đo thực chất, không phải F1. Refusal cũng thấp "
                "hơn vì ngữ cảnh sạch hơn.")

# =================================================================
# SLIDE 18 — RQ3 Embedding
# =================================================================
def slide_rq3():
    image_slide("Phần 3 · Mục 3.4", GREEN,
                "Chọn mô hình nhúng văn bản đa ngôn ngữ",
                fig("rq3_embedding_recall.png"),
                caption="Hình 12 — Recall@5/10/20 và MRR của 4 mô hình nhúng (chỉ nhánh dày đặc, 35 câu)",
                side=[
                    (0, "multilingual-e5-base ★"),
                    (1, "R@10 = 0,424 · MRR 0,277 (cao nhất)"),
                    (1, "768d · max_seq 512"),
                    (0, "e5-small: R@10 0,367 (rẻ, nhanh nhất)"),
                    (0, "mpnet thua do max_seq=128 cắt khoản dài"),
                    (0, "sbert yếu dù 768d (không InfoNCE)"),
                    (0, "→ Chọn e5-base: chất lượng cao, giữ 768d sẵn có"),
                ],
                notes_text="(~0:50 | 19:10)\n"
                "Thực nghiệm mô hình nhúng (Mục 3.4) so sánh bốn mô hình trên cùng corpus. e5-base thắng mọi chỉ số chất lượng "
                "với Recall@10 bằng 0,424. e5-small rẻ và nhanh nhất, chỉ kém khoảng 5 điểm, là phương "
                "án dự phòng tốt. multilingual-mpnet thua vì độ dài tối đa chỉ 128 token, cắt mất các "
                "khoản dài của nghị định. sbert tiếng Việt yếu dù cũng 768 chiều vì không được huấn "
                "luyện đối chiếu InfoNCE như họ E5. Em chọn e5-base vì chất lượng cao nhất và giữ "
                "nguyên 768 chiều khớp với collection đã có.")

# =================================================================
# SLIDE 19 — RQ5 Prompt
# =================================================================
def slide_rq5():
    image_slide("Phần 3 · Mục 3.6", GREEN,
                "Lệnh hướng dẫn: không HD / chỉ vai trò / đầy đủ quy tắc",
                fig("rq5_prompts_comparison.png"),
                caption="Hình 16 — Citation-P/R/F1 và tỷ lệ từ chối của 3 biến thể lệnh hướng dẫn (35 câu pháp lý)",
                side=[
                    (0, "Đầy đủ quy tắc (P2) ★"),
                    (1, "Citation-F1 = 0,279"),
                    (1, "×8,7 so không HD (0,032), ×4,2 so chỉ vai trò"),
                    (0, "Độ trễ thấp nhất 7,1s (quy tắc ép ngắn gọn)"),
                    (0, "Từ chối 0,20 = từ chối đúng lúc khi thiếu ngữ cảnh"),
                    (0, "2 biến thể đầu F1 cao hơn nhưng 'viết tuôn', không cite"),
                ],
                notes_text="(~1:00 | 20:10)\n"
                "Thực nghiệm lệnh hướng dẫn (Mục 3.6) so sánh ba biến thể. Biến thể không hướng dẫn và chỉ khai báo vai trò có F1 token cao hơn "
                "vì chúng viết tuôn ra nhiều chữ trùng đáp án, nhưng gần như không trích dẫn được. Biến "
                "thể đầy đủ 14 quy tắc (P2 trong biểu đồ) đạt Citation-F1 0,279 — gấp 8,7 lần biến thể không "
                "hướng dẫn và 4,2 lần biến thể chỉ vai trò. Thêm hai hiệu ứng tốt: độ trễ thấp nhất 7,1 giây vì quy tắc ép trả lời ngắn gọn có cấu trúc, "
                "và refusal 0,20 nghĩa là biết từ chối đúng lúc khi ngữ cảnh thiếu thay vì bịa. Một "
                "lần nữa, F1 token gây hiểu nhầm trong domain pháp lý.")

# =================================================================
# SLIDE 20 — RQ9 Query expansion
# =================================================================
def slide_rq9():
    image_slide("Phần 3 · Mục 3.7.2", GREEN,
                "Mở rộng truy vấn — mức cải thiện lớn nhất",
                fig("rq9_rewrite_comparison.png"),
                caption="Hình 19 — R@5, R@10, nDCG@10, F1 giữa không mở rộng và có mở rộng truy vấn (35 câu pháp lý)",
                side=[
                    (0, "Recall@10: 0,324 → 0,581"),
                    (1, "+26 điểm % — tăng lớn nhất toàn ablation"),
                    (0, "nDCG@10: +14 điểm"),
                    (0, "Độ trễ chỉ +7% (thêm 1 lời gọi LLM)"),
                    (0, "Đảo chiều so kho cũ: nhờ đoạn cấp Điểm nhỏ & chính xác"),
                    (0, "→ Bật mở rộng truy vấn khi vận hành thực tế"),
                ],
                notes_text="(~1:00 | 21:10)\n"
                "Thực nghiệm mở rộng truy vấn (Mục 3.7.2). Đây là kỹ thuật có tác động lớn nhất: Recall@10 tăng "
                "từ 0,324 lên 0,581, tức cộng 26 điểm phần trăm — mức tăng lớn nhất trong toàn bộ "
                "mười thực nghiệm. nDCG tăng 14 điểm, trong khi độ trễ chỉ tăng 7% do thêm một lần "
                "gọi LLM ở analyzer. Điều thú vị là kết quả này đảo chiều so với corpus phiên bản 5: "
                "khi chunk còn dài, mở rộng truy vấn tạo nhiễu; còn với corpus v6 chunk cấp Điểm nhỏ "
                "và chính xác, mở rộng truy vấn trở nên rất hữu ích. Vì vậy production bật tính năng "
                "này.")

# =================================================================
# SLIDE 21 — RQ10 Reranker + RQ4 VectorDB
# =================================================================
def slide_rq10():
    image_slide("Phần 3 · Mục 3.7.3 & 3.5", GREEN,
                "Xếp hạng lại & cơ sở dữ liệu vector — đánh đổi thực dụng",
                fig("rq10_reranker_comparison.png"),
                caption="Hình 21 — R@10, MRR, nDCG@10 & độ trễ trước/sau bge-reranker-v2-m3",
                side=[
                    (0, "Xếp hạng lại (Mục 3.7.3):"),
                    (1, "R@10 0,324 → 0,481 (+48,5%)"),
                    (1, "NHƯNG độ trễ ×156 trên CPU (~110s)"),
                    (1, "→ TẮT mặc định, bật khi có GPU"),
                    (0, "Cơ sở dữ liệu vector (Mục 3.5):"),
                    (1, "Qdrant/Chroma/FAISS R@10 ~0,42"),
                    (1, "Chênh độ trễ (4ms↔0,08ms) không đáng so LLM"),
                    (1, "→ Giữ Qdrant vì lọc payload + HTTP API"),
                ],
                notes_text="(~0:50 | 22:00)\n"
                "Thực nghiệm xếp hạng lại (Mục 3.7.3) đánh giá bộ mã hoá chéo. Nó kéo Recall@10 lên 0,481, tăng 48,5% — rất "
                "ấn tượng vì lôi được các khoản đúng nằm ở hạng 11–20 lên top 10. Nhưng cái giá là độ "
                "trễ tăng 156 lần trên CPU, khoảng 110 giây mỗi câu, không chấp nhận được cho thời "
                "gian thực. Nên reranker bị tắt mặc định, chỉ bật khi có GPU. Thực nghiệm cơ sở dữ liệu vector (Mục 3.5) so ba nền tảng: "
                "Qdrant, Chroma, FAISS recall xấp xỉ nhau; chênh lệch độ trễ vài mili giây hoàn toàn "
                "không đáng so với LLM mất vài giây. Em giữ Qdrant vì nó có payload filter để lọc "
                "status active tại tầng DB, có HTTP API và hỗ trợ cập nhật online.")

# =================================================================
# SLIDE 22 — RQ1 Architecture
# =================================================================
def slide_rq1():
    image_slide("Phần 3 · Mục 3.8", GREEN,
                "So sánh kiến trúc hệ thống & 'bẫy' của F1",
                plot("rq1_pipeline_comparison.png"),
                caption="Hình 22 — F1, ROUGE-L, độ trễ, độ chính xác phân loại: chỉ Gemini / RAG cơ bản / RAG hướng tác tử",
                side=[
                    (0, "RAG hướng tác tử ★ Phân loại = 0,975"),
                    (1, "vs RAG cơ bản 0,875 · định tuyến 100% câu OOS đúng"),
                    (0, "Lưu ý: chỉ Gemini F1 cao nhất (0,370) là SAI LỆCH:"),
                    (1, "câu OOS từ chối trùng n-gram + bịa nguồn dài"),
                    (0, "→ F1 token gây hiểu nhầm; phân loại & Cit-R mới thực chất"),
                    (0, "Độ trễ tác tử bị giãn cách 2s; thực tế ~20–25s"),
                ],
                notes_text="(~1:10 | 23:10)\n"
                "Thực nghiệm so sánh kiến trúc (Mục 3.8) trên đủ 40 câu. RAG hướng tác tử đạt độ chính xác phân loại "
                "0,975 — phân loại đúng 39/40 câu, định tuyến 100% câu ngoài phạm vi chính xác — so "
                "với RAG cơ bản chỉ 0,875. Điểm cần giải thích kỹ: kiến trúc chỉ dùng Gemini lại có F1 cao nhất 0,370, "
                "nhưng đây là một artifact của metric. Thứ nhất, với câu ngoài phạm vi, Gemini từ "
                "chối tự nhiên, n-gram trùng nhiều với đáp án vốn cũng là câu từ chối. Thứ hai, "
                "Gemini bịa nguồn pháp lý dài chi tiết, kéo F1 lên nhờ trùng từ bề mặt nhưng trích "
                "dẫn sai hoàn toàn. Đây đúng là hiện tượng ảo giác pháp lý. Vì vậy F1 token gây hiểu "
                "nhầm; Category Accuracy và Citation-Recall mới phản ánh chất lượng thật. Độ trễ "
                "agentic 52 giây bị thổi phồng do throttle 2 giây mỗi call; thực tế không throttle "
                "khoảng 20–25 giây.")

# =================================================================
# SLIDE 23 — RQ7 RAGAS
# =================================================================
def slide_rq7():
    image_slide("Phần 3 · Mục 3.9", GREEN,
                "Đánh giá chất lượng câu trả lời bằng RAGAS (gpt-4o-mini)",
                fig("rq7_ragas.png"),
                caption="Hình 23 — Phân phối độ chính xác ngữ cảnh, độ trung thực, độ phủ ngữ cảnh (n=35)",
                side=[
                    (0, "context_precision = 0,885 (rất cao)"),
                    (1, "chunk lấy lên đa số liên quan, ít nhiễu"),
                    (0, "faithfulness = 0,550"),
                    (1, "hơn nửa khẳng định bám ngữ cảnh"),
                    (0, "context_recall = 0,452 = BOTTLENECK"),
                    (1, "mới bắt ~45% facts cần thiết"),
                    (0, "answer_relevancy: lỗi lib RAGAS 0.2 (không kết luận)"),
                ],
                notes_text="(~1:00 | 24:10)\n"
                "Thực nghiệm đánh giá chất lượng câu trả lời (Mục 3.9) dùng RAGAS với gpt-4o-mini làm giám khảo độc lập. Độ chính xác ngữ cảnh đạt 0,885 "
                "rất cao, xác nhận các chunk đưa vào ngữ cảnh đa số thực sự liên quan, ít nhiễu — nhờ "
                "truy xuất lai cộng sibling enrichment. Faithfulness 0,550, tức hơn một nửa khẳng định "
                "trong câu trả lời được ngữ cảnh hỗ trợ, tốt hơn hẳn báo cáo cũ. Nhưng context recall "
                "chỉ 0,452 — đây mới là bottleneck thật của hệ thống: retrieval mới bắt được khoảng "
                "45% thông tin cần thiết. Chính vì vậy reranker và mở rộng truy vấn được ưu tiên "
                "trong roadmap. Riêng answer_relevancy bị lỗi thư viện RAGAS 0.2 nên không kết luận.")

# =================================================================
# SLIDE 24 — RQ8 Cost/Perf
# =================================================================
def slide_rq8():
    image_slide("Phần 3 · Mục 3.10", GREEN,
                "Phân tích chi phí & hiệu năng vận hành",
                fig("rq8_node_breakdown.png"),
                caption="Hình 24 — Phân bố chi phí & thời gian theo từng nút (truy vết, 1.964 lượt chạy)",
                side=[
                    (0, "Cost = $0,00344 / câu"),
                    (1, "legal_rag chiếm 77%"),
                    (1, "≈ $46/tháng cho 10.000 query"),
                    (0, "Latency agentic mean 33s (có throttle); ~20–25s thực"),
                    (0, "Gemini error rate = 14,6%"),
                    (1, "free-tier rate limit; cần stable paid tier → <2%"),
                ],
                notes_text="(~0:50 | 25:00)\n"
                "Thực nghiệm chi phí & hiệu năng (Mục 3.10) phân tích từ gần hai nghìn lượt chạy trên hệ thống truy vết. Mỗi câu tốn "
                "khoảng 0,0034 đô la, trong đó nút legal_rag chiếm 77%; ước tính mười nghìn truy vấn "
                "mỗi tháng chỉ khoảng 46 đô la tiền LLM. Độ trễ trung bình 33 giây có throttle, thực "
                "tế 20–25 giây. Điểm cần lưu ý là tỉ lệ lỗi Gemini 14,6% do dùng bản preview free "
                "tier bị giới hạn tần suất; chuyển sang bản stable trả phí sẽ giảm xuống dưới 2%.")

# =================================================================
# SLIDE 25 — Bảng tổng hợp RQ
# =================================================================
def slide_rq_summary():
    rows = [
        ["Phân đoạn", "Phân cấp Điều→Khoản→Điểm", "3.3", "Cit-R 0,41 vs 0,17 (×2,4)"],
        ["Mô hình nhúng", "multilingual-e5-base 768d", "3.4", "R@10 0,42 (cao nhất 4 MH)"],
        ["Cơ sở DL vector", "Qdrant 1.17 HNSW + chỉ mục payload", "3.5", "R@10 0,42, lọc payload"],
        ["Truy xuất thưa", "BM25 Okapi (k₁=1,5; b=0,75) + pyvi", "3.7.1", "Truy xuất lai R@10 0,324"],
        ["Mở rộng truy vấn", "BẬT", "3.7.2", "R@10 0,581 vs 0,324 (+26 pts)"],
        ["Xếp hạng lại", "TẮT mặc định (bật khi có GPU)", "3.7.3", "+48,5% R@10 nhưng ×156 độ trễ"],
        ["Lệnh hướng dẫn", "Đầy đủ quy tắc + lọc trích dẫn", "3.6", "Cit-F1 0,279 vs 0,032 (×8,7)"],
        ["Kiến trúc", "RAG hướng tác tử (LangGraph 6 nút)", "3.8", "Phân loại 0,975 vs 0,875"],
        ["LLM", "Gemini 3.1 Flash Lite Preview", "3.10", "0,00344 USD/câu · lỗi 14,6%"],
        ["Giám khảo", "RAGAS + gpt-4o-mini", "3.9", "chính xác 0,885 · phủ 0,452"],
    ]
    table_slide("Phần 3 · Tổng hợp", GREEN,
                "Tổng hợp quyết định thiết kế từ thực nghiệm (Bảng 25)",
                ["Thành phần", "Lựa chọn", "Mục", "Số liệu chính"], rows,
                col_w=[2.3, 4.7, 1.0, 4.3], font_size=13, hl_rows=[1, 5, 7, 8],
                notes_text="(~0:20 | 25:20)\n"
                "Bảng này tổng hợp toàn bộ mười thực nghiệm: mỗi thành phần thiết kế đều có bằng chứng "
                "định lượng đi kèm, kèm số mục tương ứng trong báo cáo. Bốn dòng được tô đậm là các "
                "quyết định có tác động lớn nhất: phân đoạn phân cấp, mở rộng truy vấn, lệnh hướng dẫn "
                "đầy đủ quy tắc, và kiến trúc hướng tác tử.")

# =================================================================
# SLIDE 26 — Triển khai production
# =================================================================
def slide_deploy():
    s = base("Phần 4 · Triển khai", PURPLE)
    title_on(s, "Triển khai production — $0/tháng & CI/CD")
    # sơ đồ deploy
    box(s, 0.6, 2.0, 1.8, 0.8, "Trình duyệt\nngười dùng", TEAL, WHITE, 12, True)
    fe = box(s, 3.0, 2.0, 2.7, 0.85, "Vercel\nNext.js 14 (FE)\nNextAuth", GREEN, WHITE, 12, True)
    be = box(s, 6.3, 2.0, 2.9, 0.85, "HF Spaces\nFastAPI + LangGraph\n2 vCPU · 16GB", BLUE, WHITE, 12, True)
    qd = box(s, 9.9, 1.45, 2.9, 0.8, "Qdrant Cloud\n4.423 × 768d", PINK, WHITE, 12, True)
    pg = box(s, 9.9, 2.5, 2.9, 0.8, "Supabase Postgres\nAuth + checkpoint", PINK, WHITE, 12, True)
    arrow(s, 2.4, 2.4, 3.0, 2.4, GREY, 1.5)
    arrow(s, 5.7, 2.4, 6.3, 2.4, GREY, 1.5)
    arrow(s, 9.2, 2.25, 9.9, 1.85, GREY, 1.3)
    arrow(s, 9.2, 2.6, 9.9, 2.9, GREY, 1.3)
    arrow(s, 4.35, 2.85, 10.5, 3.3, GREY, 1.0, dashed=True)  # FE -> Supabase (Prisma)
    box(s, 6.3, 3.45, 2.9, 0.5, "Gemini · Tavily (dịch vụ ngoài)", PURPLE, WHITE, 11, True)
    arrow(s, 7.7, 2.85, 7.7, 3.45, GREY, 1.1, dashed=True)
    box(s, 0.6, 4.3, 5.9, 0.46, "Hạ tầng $0/tháng (free tier)", PURPLE, WHITE, 14, True)
    bullets_box(s, 0.7, 4.85, 5.8, 2.0, [
        (0, "Cùng 1 codebase & Docker image cho local + cloud"),
        (1, "Khác biệt chỉ ở biến môi trường"),
        (0, "Local: Docker Compose (Qdrant + FastAPI + Next.js)"),
        (0, "Observability: /metrics, /pending, LangSmith always-on"),
    ], base_size=14, gap=6)
    box(s, 6.7, 4.3, 6.1, 0.46, "CI/CD — cổng hồi quy tự động", ORANGE, WHITE, 14, True)
    bullets_box(s, 6.8, 4.85, 6.0, 2.0, [
        (0, "GitHub Actions: job smoke (6 unit test) + job regression"),
        (0, "Assert: mean Recall@10 ≥ 0,40 trên 35 câu legal"),
        (0, "Không chunk status=repealed lọt kết quả · latency < 2s"),
        (1, "Fail → block PR (chống cập nhật corpus làm tụt chất lượng)"),
    ], base_size=14, gap=6)
    notes(s, "(~1:20 | 26:40)\n"
             "Hệ thống đã triển khai thực tế trên đám mây miễn phí, tổng chi phí hạ tầng bằng 0. "
             "Frontend Next.js chạy trên Vercel, backend FastAPI cộng LangGraph trên Hugging Face "
             "Spaces, kho vector trên Qdrant Cloud, còn xác thực và checkpoint dùng chung một "
             "Supabase Postgres; Gemini và Tavily là dịch vụ ngoài. Điểm hay là cùng một codebase và "
             "Docker image phục vụ cả local lẫn cloud, chỉ khác biến môi trường. Về chất lượng, em "
             "thiết lập CI/CD trên GitHub Actions với hai job: smoke test và regression test. Job "
             "regression khởi động Qdrant, nạp 4.423 chunk và bắt buộc Recall@10 trung bình phải từ "
             "0,40 trở lên, không được có chunk đã bãi bỏ lọt vào, độ trễ dưới 2 giây — nếu fail thì "
             "chặn pull request. Cơ chế này ngăn việc cập nhật corpus vô tình làm tụt chất lượng truy "
             "xuất.")

# =================================================================
# SLIDE 27 — Giao diện sản phẩm
# =================================================================
def slide_ui():
    two_col_slide("Phần 4 · Sản phẩm", PURPLE,
                  "Giao diện sản phẩm",
                  "Chat người dùng", [
                      (0, "Streaming real-time từng token (SSE)"),
                      (0, "Citation panel: 'Điều X, Khoản Y, Điểm Z — [doc_id]'"),
                      (0, "Hội thoại đa lượt (checkpointer theo thread_id)"),
                      (1, "'Còn xe tải hạng nặng thì sao?'"),
                      (0, "Nhãn cảnh báo 'Tra cứu từ internet' cho câu web"),
                      (0, "Trạng thái pending khi chờ duyệt HITL"),
                  ],
                  "Admin Console (/admin) — HITL", [
                      (0, "Chỉ tài khoản trong ADMIN_EMAILS"),
                      (0, "Liệt kê thread pending từ /pending"),
                      (0, "Xem: câu hỏi gốc + 3 snippet Tavily + draft"),
                      (0, "Approve → resume web_finalize"),
                      (0, "Reject → xoá thread, trả thông báo từ chối"),
                      (0, "Timeout 10 phút; state vẫn lưu trong DB"),
                  ],
                  lc=TEAL, rc=RED,
                  notes_text="(~0:50 | 27:30)\n"
                  "Về sản phẩm, giao diện chat xây trên Next.js 14: câu trả lời hiện dần từng token "
                  "qua SSE, có panel trích dẫn hiển thị chính xác điều khoản điểm và mã văn bản để "
                  "người dùng tự kiểm chứng, hỗ trợ hội thoại nhiều lượt nên câu sau có thể tham "
                  "chiếu câu trước, và gắn nhãn cảnh báo cho câu trả lời từ internet. Bên cạnh là "
                  "trang quản trị cho con người trong vòng lặp: quản trị viên thấy danh sách thread "
                  "đang chờ, xem câu hỏi và các snippet Tavily rồi bấm Approve hoặc Reject. Hai ảnh "
                  "chụp màn hình thực tế (Hình 28 — giao diện chat; Hình 29 — bảng quản trị HITL) "
                  "em sẽ chèn ở slide này khi đóng quyển.")

# =================================================================
# SLIDE 28 — Demo minh hoạ
# =================================================================
def slide_demo():
    s = base("Phần 4 · Minh hoạ", PURPLE)
    title_on(s, "Minh hoạ thực tế — đối chiếu hệ thống vs đáp án chuẩn")
    box(s, 0.5, 1.7, 12.3, 0.5, "Câu hỏi: \"Vượt đèn đỏ đi xe máy phạt bao nhiêu tiền?\"  (danh mục phạt đơn nghĩa)",
        NAVY, WHITE, 14, True)
    box(s, 0.5, 2.35, 6.0, 0.42, "Đáp án chuẩn", GREY, WHITE, 13, True)
    bullets_box(s, 0.6, 2.85, 5.9, 1.7, [
        (0, "Phạt 800.000–1.000.000đ + trừ 2 điểm GPLX"),
        (1, "[168/2024/NĐ-CP · Điều 7 · Khoản 4 · điểm p]"),
    ], base_size=14, gap=6)
    box(s, 6.83, 2.35, 6.0, 0.42, "Đầu ra của hệ thống", GREEN, WHITE, 13, True)
    bullets_box(s, 6.93, 2.85, 5.9, 1.7, [
        (0, "Phạt 800.000–1.000.000đ [Đ7·K4·p] ✓"),
        (0, "+ Làm giàu lân cận kéo Khoản 11: trừ 2 điểm [Đ7·K11] ✓"),
        (1, "→ Đầy đủ hơn cả đáp án chuẩn tối thiểu"),
    ], base_size=14, gap=6)
    box(s, 0.5, 4.75, 12.3, 0.5,
        "Câu đa ý định: \"Vượt đèn đỏ + không mang bằng lái xe máy phạt tổng bao nhiêu?\"  (danh mục đa ý định)",
        NAVY, WHITE, 13, True)
    bullets_box(s, 0.6, 5.35, 12.2, 1.4, [
        (0, "Hệ thống tách đúng 2 ý: Đ7·K4·p (800K–1tr) + Đ24·K2·c (100–200K) → tổng 900K–1,2tr"),
        (0, "Khớp đầy đủ 2 trích dẫn chuẩn · in đậm số tiền · cộng tổng (lệnh hướng dẫn quy tắc 6)"),
    ], base_size=15, gap=8)
    notes(s, "(~1:00 | 28:30)\n"
             "Em minh hoạ bằng hai ví dụ thực từ bộ đánh giá. Câu đơn giản 'vượt đèn đỏ đi xe máy "
             "phạt bao nhiêu': hệ thống trả đúng mức 800 nghìn đến 1 triệu, trích dẫn đúng Điều 7 "
             "Khoản 4 điểm p; hơn nữa nhờ sibling enrichment nó còn kéo thêm Khoản 11 về trừ 2 điểm "
             "giấy phép — tức đầy đủ hơn cả đáp án chuẩn tối thiểu. Câu đa ý định 'vượt đèn đỏ và "
             "không mang bằng lái': hệ thống tách đúng hai hành vi, trích dẫn đúng hai điều khoản, in "
             "đậm số tiền và cộng tổng — đúng như quy tắc 6 của lệnh hướng dẫn. Nếu thầy cô cho phép, em "
             "có thể demo trực tiếp trên hệ thống đang chạy.")

# =================================================================
# SLIDE 29 — Kết luận
# =================================================================
def slide_conclusion():
    s = base("Phần 4 · Kết luận", PURPLE)
    title_on(s, "Kết luận — kết quả · hạn chế · hướng phát triển")
    box(s, 0.5, 1.7, 4.0, 0.44, "Kết quả đạt được", GREEN, WHITE, 13, True)
    bullets_box(s, 0.55, 2.2, 3.95, 4.4, [
        (0, "R@10 = 0,581 (×2,9 so baseline 0,200)"),
        (0, "Citation-F1 = 0,279 (×8,7 zero-shot)"),
        (0, "Category Acc = 0,975"),
        (0, "context_precision = 0,885"),
        (0, "$0,00344/câu · ~$46/10k query"),
        (0, "Production $0/tháng + CI gate"),
    ], base_size=13, gap=8)
    box(s, 4.67, 1.7, 4.0, 0.44, "Hạn chế", RED, WHITE, 13, True)
    bullets_box(s, 4.72, 2.2, 3.95, 4.4, [
        (0, "Eval set nhỏ (35 câu legal)"),
        (0, "Gemini free-tier error 14,6%"),
        (0, "Reranker chưa dùng được prod (cần GPU)"),
        (0, "context_recall 0,452 = bottleneck"),
        (0, "answer_relevancy lỗi lib"),
        (0, "Refusal 20% còn false-negative"),
    ], base_size=13, gap=8)
    box(s, 8.84, 1.7, 4.0, 0.44, "Hướng phát triển (Bảng 27)", PURPLE, WHITE, 13, True)
    bullets_box(s, 8.89, 2.2, 3.95, 4.4, [
        (0, "Ngắn hạn:"),
        (1, "Mở rộng eval 80–100 câu"),
        (1, "Gemini stable trả phí → lỗi <2%"),
        (1, "Tối ưu chi phí; củng cố vận hành"),
        (0, "Trung hạn:"),
        (1, "Bật reranker GPU (~2–3s)"),
        (1, "Graph RAG trên đồ thị tri thức"),
        (1, "Đánh giá lại bằng giám khảo mạnh hơn"),
        (0, "Dài hạn:"),
        (1, "Kiến trúc đa tác tử; truy xuất đa bước"),
    ], base_size=13, gap=5)
    notes(s, "(~1:00 | 29:30)\n"
             "Tổng kết, hệ thống đạt Recall@10 0,581 — gấp 2,9 lần phân đoạn cố định; Citation-F1 0,279 — gấp "
             "8,7 lần không hướng dẫn; độ chính xác phân loại 0,975; độ chính xác ngữ cảnh 0,885; chi phí chỉ 0,0034 "
             "đô mỗi câu và đã chạy thực tế miễn phí với cổng CI bảo vệ chất lượng. Về hạn chế: bộ "
             "đánh giá còn nhỏ, Gemini free tier lỗi 14,6%, reranker chưa dùng được vì cần GPU, "
             "độ phủ ngữ cảnh 0,452 vẫn là điểm nghẽn, và độ liên quan câu trả lời bị lỗi thư viện. Hướng phát "
             "triển chia ba tầm: ngắn hạn mở rộng bộ đánh giá lên 80–100 câu, chuyển Gemini sang bản "
             "stable trả phí, tối ưu chi phí và củng cố vận hành; trung hạn bật reranker khi có GPU, "
             "xây Graph RAG trên đồ thị tri thức pháp luật để nâng độ phủ ngữ cảnh, và đánh giá lại "
             "bằng giám khảo mạnh hơn; dài hạn tiến tới kiến trúc đa tác tử với tác tử truy xuất đa bước.")

# =================================================================
# SLIDE 30 — Cảm ơn
# =================================================================
def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
    textbox(s, 0.8, 2.5, 11.7, 1.0, "Cảm ơn quý thầy cô đã lắng nghe",
            34, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, 0.8, 3.7, 11.7, 0.6, "Em xin sẵn sàng trả lời các câu hỏi của hội đồng",
            18, RGBColor(0xC5, 0xCA, 0xE9), False, PP_ALIGN.CENTER)
    box(s, 3.9, 4.7, 5.5, 0.7, "Mạc Phú Phong  ·  106210059", ORANGE, WHITE, 16, True)
    notes(s, "(~0:20 | 29:50)\n"
             "Phần trình bày của em đến đây là hết. Em cảm ơn quý thầy cô đã lắng nghe và rất mong "
             "nhận được các câu hỏi, góp ý từ hội đồng.")

# ---------- Dựng theo thứ tự ----------
slide_title()
slide_agenda()
slide_problem()
slide_goals()
slide_rag()
slide_foundations()
slide_agentic()
slide_arch()
slide_stack()
slide_clean()
slide_chunk()
slide_index()
slide_retrieval()
slide_prompt()
slide_graph()
slide_exp_design()
slide_rq2()
slide_rq3()
slide_rq5()
slide_rq9()
slide_rq10()
slide_rq1()
slide_rq7()
slide_rq8()
slide_rq_summary()
slide_deploy()
slide_ui()
slide_demo()
slide_conclusion()
slide_thanks()

prs.save(OUT)
print("Saved:", OUT)
print("Tổng số slide:", len(prs.slides._sldIdLst))
